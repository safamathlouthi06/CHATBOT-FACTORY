"""
Extraction texte PDF & PPT
"""

from PyPDF2 import PdfReader
from pptx import Presentation


def extract_text_from_pdf(file_path: str) -> str:
    from PyPDF2 import PdfReader

    reader = PdfReader(file_path)
    text = ""

    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"

    # ✅ nettoyage IMPORTANT
    text = text.replace("\n", " ")
    text = text.replace("  ", " ")

    return text.strip()



def extract_text_from_ppt(file_path: str) -> str:
    text = ""
    prs = Presentation(file_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text

